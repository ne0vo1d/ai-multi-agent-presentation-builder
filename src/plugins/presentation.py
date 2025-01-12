from typing import Annotated
from semantic_kernel.functions.kernel_function_decorator import kernel_function
from pptx import Presentation


class PresentationPlugin:
    """The Presentation Plugin can be used to create presentations decks and slides."""

    @kernel_function(description="Create a presentation deck in a PDF format.")
    def create_presentation(self, 
                            title: Annotated[str, "The title of the presentation"],
                            subtitle: Annotated[str, "The subtitle of the presentation"],
                            content: Annotated[str, "The content of the decks"]) -> Annotated[str, "Create a presentation."]:

        # Create a presentation object
        prs = Presentation(pptx='green.pptx')

        # Remove title and subtitle from the slides content
        content = content.replace(title, '')
        content = content.replace(subtitle, '')

        # Split the content into slides
        slides_content = content.split("#")
        initial_slide = True
        
        for slide_content in slides_content:
            # The first slide must have a title and subtitle
            if initial_slide:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                slide.placeholders[1].text = subtitle
                initial_slide = False
                continue
            
            if len(slide_content) > 0:
                title = slide_content.split('\n')[0]
                slide_content = slide_content.replace(title, '').strip()
                # Add the content for the other slides
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                slide.placeholders[1].text = slide_content

        # Save the presentation
        output_path = "presentation.pptx"
        prs.save(output_path)

        return output_path
